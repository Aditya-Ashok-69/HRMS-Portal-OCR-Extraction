# document_processor.py
import os
import re
from pathlib import Path

from extractors.aadhaar import extract_aadhaar
from extractors.pan import extract_pan
from extractors.llm_extractor import extract_with_llm, SCHEMAS

MAX_FILE_SIZE_BYTES = 10 * 1024 * 1024  # 10MB

ALLOWED_EXTENSIONS = {
    "aadhaar":            {".pdf"},
    "pan":                {".jpg", ".jpeg", ".png"},
    "resume":             {".pdf", ".docx", ".pptx"},
    "payslip":            {".pdf"},
    "experience_letter":  {".pdf", ".docx"},
    "degree_certificate": {".pdf", ".jpg", ".jpeg", ".png"},
}


def process_document(file_path, doc_type: str, sub_type: str = "image", **kwargs) -> dict:
    size = os.path.getsize(file_path)
    if size > MAX_FILE_SIZE_BYTES:
        return {"error": f"File too large: {size // (1024*1024)}MB. Max 10MB.", "doc_type": doc_type}

    ext = Path(file_path).suffix.lower()
    allowed = ALLOWED_EXTENSIONS.get(doc_type, set())
    if allowed and ext not in allowed:
        return {"error": f"Invalid file type '{ext}' for {doc_type}. Allowed: {sorted(allowed)}", "doc_type": doc_type}

    if doc_type == "aadhaar":
        result = extract_aadhaar(file_path, sub_type="pdf")
        #print("DEBUG aadhaar:", type(result), result)
        return _to_output_json(result)

    if doc_type == "pan":
        result = extract_pan(file_path, sub_type="image")
        #print("DEBUG pan:", type(result), result)
        return _to_output_json(result)

    if doc_type in SCHEMAS or doc_type == "resume":
        raw_text = _get_text(file_path)
        result = extract_with_llm(raw_text, doc_type, **kwargs)
        #print(f"DEBUG {doc_type}:", type(result), result)
        return _to_output_json(result)

    return {"error": f"Unsupported doc_type: {doc_type}"}


def _to_output_json(result) -> dict:
    exclude = {"raw_text", "_meta"}

    # If extractor returned None or a non-dict, surface that as an error
    if not isinstance(result, dict):
        return {
            "error": "internal_error: extractor returned no result",
            "raw_result": str(result)
        }

    return {k: v for k, v in result.items() if k not in exclude}


def _get_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        from extract_id import extract_text_from_pdf
        import pypdfium2 as pdfium
        from extract_id import paddle_ocr_image, preprocess_for_ocr

        text = extract_text_from_pdf(file_path)
        if len(text.strip()) < 200 or not re.search(r'\d{4,}', text):
            # Scanned PDF — force OCR
            pdf = pdfium.PdfDocument(file_path)
            pages = []
            try:
                for page in pdf:
                    img = page.render(scale=3).to_pil()
                    pages.append(paddle_ocr_image(preprocess_for_ocr(img)))
            finally:
                pdf.close()
            return "\n".join(pages)
        return text

    elif ext in {".jpg", ".jpeg", ".png", ".tiff"}:
        from extract_id import paddle_ocr_image
        from PIL import Image
        return paddle_ocr_image(Image.open(file_path))

    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return ""  # caller will get empty extraction

    elif ext in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
            return "\n".join(lines)
        except ImportError:
            return ""

    return ""

